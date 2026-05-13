import json
import base64
import io
import pymupdf
import os
from pptx import Presentation
from PIL import Image
from docx import Document
import pdfplumber
import asyncio
import logging
import time
import pytesseract
from pdf2image import convert_from_bytes
import psutil
import gc
import tempfile

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

from .qdrant_setup import *
from .agents.generation_agent import generation_agent, generation_agent_async
from .agents.schema_agent import generate_dataset_schema
from .agents.evolution_agent.evolver import evolve_dataset
from .utils import process_datagen_prompt

# Memory threshold for hybrid processing (using existing tier limits)
MEMORY_THRESHOLD_FREE = 5 * 1024 * 1024  # 5MB for free tier
MEMORY_THRESHOLD_PAID = 50 * 1024 * 1024  # 50MB for paid tier
def process_document(path):
    # Open the PDF using the single, consistent namespace
    pdf = pymupdf.open(path)
    return pdf

def get_memory_threshold(user_tier: str = "free") -> int:
    """Get memory threshold based on user tier"""
    if not user_tier:
        logger.warning("⚠️ No user tier provided, defaulting to free tier")
        user_tier = "free"
    
    valid_tiers = ['free', 'paid', 'enterprise']
    if user_tier not in valid_tiers:
        logger.warning(f"⚠️ Invalid user tier '{user_tier}', defaulting to free tier")
        user_tier = "free"
    return MEMORY_THRESHOLD_PAID if user_tier in ['paid', 'enterprise'] else MEMORY_THRESHOLD_FREE

async def process_document_hybrid(file_bytes: bytes, filename: str, user_tier: str = "free") -> list:
    """
    Process document using hybrid approach: small files in memory, large files on disk
    """
    file_size = len(file_bytes)
    memory_threshold = get_memory_threshold(user_tier)
    
    logger.info(f"📁 Processing {filename} ({file_size / 1024 / 1024:.2f} MB) with {user_tier} tier")
    
    if file_size <= memory_threshold:
        # Small file: process entirely in memory
        logger.info(f"🧠 Processing {filename} in memory (≤{memory_threshold / 1024 / 1024:.1f}MB)")
        return await process_document_in_memory(file_bytes, filename)
    else:
        # Large file: use disk processing (current approach)
        logger.info(f"💾 Processing {filename} on disk (>{memory_threshold / 1024 / 1024:.1f}MB)")
        return await process_document_on_disk(file_bytes, filename)

async def process_document_in_memory(file_bytes: bytes, filename: str) -> list:
    """
    Process document entirely in memory with comprehensive logging and cleanup
    """
    process = psutil.Process()
    initial_memory = process.memory_info().rss / 1024 / 1024  # MB
    after_memory = None
    logger.info(f"📊 Memory before processing {filename}: {initial_memory:.2f} MB")
    
    try:
        # Convert to PDF if necessary
        converted_pdf_bytes = convert_to_pdf(file_bytes, filename)
        logger.info(f"🔄 Converted {filename} to PDF format")
        
        # Extract text and create chunks
        chunks = await create_hybrid_chunks_in_memory(converted_pdf_bytes, filename)
        
        # Log memory after processing
        after_memory = process.memory_info().rss / 1024 / 1024
        logger.info(f"📊 Memory after processing {filename}: {after_memory:.2f} MB")
        logger.info(f"📊 Memory used: {after_memory - initial_memory:.2f} MB")
        
        return chunks
        
    except Exception as e:
        logger.error(f"❌ Error processing {filename} in memory: {e}")
        raise
    finally:
        # Clear file bytes from memory
        del file_bytes
        if 'converted_pdf_bytes' in locals():
            del converted_pdf_bytes
        
        # Force garbage collection
        gc.collect()
        
        # Log final memory
        if after_memory is not None:  # ← Check if defined
            final_memory = process.memory_info().rss / 1024 / 1024
            logger.info(f"📊 Memory freed: {after_memory - final_memory:.2f} MB")
            logger.info(f"📊 Memory after cleanup {filename}: {final_memory:.2f} MB")
        else:
            logger.info("📊 Memory after cleanup was not measured")
async def process_document_on_disk(file_bytes: bytes, filename: str) -> list:
    """
    Process document on disk (fallback for large files)
    """
    logger.info(f"💾 Processing {filename} on disk")
    
    # Create temporary directory
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = os.path.join(temp_dir, filename)
        
        # Write file to temp directory
        with open(temp_path, 'wb') as f:
            f.write(file_bytes)
        
        # Process using existing disk-based approach
        chunks = create_chunks_from_disk(temp_path, filename)
        
        return chunks

async def create_hybrid_chunks_in_memory(pdf_bytes: bytes, filename: str, max_chunk_size: int = 4000) -> list:
    """
    Create chunks in memory using hybrid approach: respect page boundaries but optimize large pages
    """
    logger.info(f"🧠 Creating hybrid chunks for {filename} in memory")
    
    # Extract text from PDF
    pages = extract_text_from_pdf(pdf_bytes)
    chunks = []
    
    for page_num, page_content in enumerate(pages):
        if len(page_content) <= max_chunk_size:
            # Small page: keep as single chunk
            chunks.append({
                "filename": filename,
                "page_number": page_num + 1,
                "page_content": page_content,
                "chunk_type": "full_page",
                "chunk_size": len(page_content)
            })
            logger.debug(f"📄 Page {page_num + 1}: Single chunk ({len(page_content)} chars)")
        else:
            # Large page: split by sentences while respecting page boundaries
            page_chunks = split_page_by_sentences(page_content, max_chunk_size, page_num + 1)
            chunks.extend(page_chunks)
            logger.debug(f"📄 Page {page_num + 1}: Split into {len(page_chunks)} chunks")
    
    logger.info(f"✅ Created {len(chunks)} chunks for {filename}")
    return chunks

def split_page_by_sentences(page_content: str, max_chunk_size: int, page_number: int) -> list:
    """
    Split a large page into semantic chunks while respecting sentence boundaries
    """
    chunks = []
    
    # Split by sentences (simple approach - can be enhanced later)
    sentences = page_content.split('. ')
    
    current_chunk = ""
    chunk_part = 1
    
    for sentence in sentences:
        # Add period back to sentence
        sentence_with_period = sentence + ". "
        
        if len(current_chunk) + len(sentence_with_period) <= max_chunk_size:
            # Add sentence to current chunk
            current_chunk += sentence_with_period
        else:
            # Save current chunk if it has content
            if current_chunk.strip():
                chunks.append({
                    "filename": f"page_{page_number}_part_{chunk_part}",
                    "page_number": page_number,
                    "chunk_part": chunk_part,
                    "content": current_chunk.strip(),
                    "chunk_type": "page_fragment",
                    "chunk_size": len(current_chunk.strip())
                })
                chunk_part += 1
            
            # Start new chunk with current sentence
            current_chunk = sentence_with_period
    
    # Add the last chunk if it has content
    if current_chunk.strip():
        chunks.append({
            "filename": f"page_{page_number}_part_{chunk_part}",
            "page_number": page_number,
            "chunk_part": chunk_part,
            "content": current_chunk.strip(),
            "chunk_type": "page_fragment",
            "chunk_size": len(current_chunk.strip())
        })
    
    return chunks

def create_chunks_from_disk(directory_path: str, filename: str = None) -> list:
    """
    Legacy disk-based chunking (fallback for large files)
    """
    logger.info(f"💾 Using disk-based chunking for {filename or directory_path}")
    
    if os.path.isfile(directory_path):
        # Single file
        with open(directory_path, "rb") as f:
            file_bytes = f.read()
        
        filename = filename or os.path.basename(directory_path)
        converted_pdf_bytes = convert_to_pdf(file_bytes, filename)
        pages = extract_text_from_pdf(converted_pdf_bytes)
        
        chunks = []
        for page_num, page in enumerate(pages):
            chunks.append({
                "filename": filename, 
                "page_number": page_num + 1,
                "page_content": page,
                "chunk_type": "legacy_disk",
                "chunk_size": len(page)
            })
        return chunks
    
    else:
        # Directory (multiple files)
        file_paths = [
            os.path.abspath(os.path.join(directory_path, f))
            for f in os.listdir(directory_path)
            if os.path.isfile(os.path.join(directory_path, f))
        ]
        
        all_chunks = []
        for file_path in file_paths:
            filename = os.path.basename(file_path)
            with open(file_path, "rb") as f:
                file_bytes = f.read()
            
            converted_pdf_bytes = convert_to_pdf(file_bytes, filename)
            pages = extract_text_from_pdf(converted_pdf_bytes)
            
            for page_num, page in enumerate(pages):
                all_chunks.append({
                    "filename": filename, 
                    "page_number": page_num + 1,
                    "page_content": page,
                    "chunk_type": "legacy_disk",
                    "chunk_size": len(page)
                })
        
        return all_chunks

# Legacy function for backward compatibility
def create_chunks(directory_path: str):
    """
    Legacy chunking function - now uses hybrid approach
    """
    logger.warning("⚠️ Using legacy create_chunks function. Consider using process_document_hybrid instead.")
    return create_chunks_from_disk(directory_path)

def encode_pdf(pdf_bytes: bytes):
    """Encode PDF bytes to a base64 string."""
    try:
        return base64.b64encode(pdf_bytes).decode("utf-8")
    except Exception as e:
        print(f"Error encoding PDF to base64: {e}")
        return None

def convert_to_pdf(file_bytes: bytes, filename: str):
    extension = filename.lower().split('.')[-1]

    if extension == "pdf":
        return file_bytes

    buffer = io.BytesIO()
    pdf = pymupdf.open()

    if extension in {"jpg", "jpeg", "png", "gif", "webp", "bmp"}:
        img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
        img.save(buffer, format="PDF")
        return buffer.getvalue()

    elif extension in {"txt", "md"}:
        text = file_bytes.decode("utf-8", errors="ignore")
        page = pdf.new_page()
        page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    elif extension in {"doc", "docx"}:
        doc = Document(io.BytesIO(file_bytes))
        text = "\n".join([para.text for para in doc.paragraphs])
        page = pdf.new_page()
        page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    elif extension == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            page = pdf.new_page()
            page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    else:
        raise ValueError(f"Unsupported file type: {extension}")
    
def process_page(idx, ocr_response=None):
    try:
        if ocr_response and hasattr(ocr_response, 'pages') and idx < len(ocr_response.pages):
            return ocr_response.pages[idx].markdown
        else:
            return f"Error: Page {idx + 1} not available in OCR response"
    except Exception as e:
        return f"Error processing page {idx + 1}: {e}"

def extract_text_from_pdf(pdf_bytes: bytes):
    logger.info("🔍 Starting PDF text extraction...")
    start_time = time.time()
    last_error = None

    # First try traditional PDF text extraction with pdfplumber
    logger.info("📖 Attempting text extraction with pdfplumber...")
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            extracted_text = []
            for page_num, page in enumerate(pdf.pages):
                logger.debug(f"Processing page {page_num + 1} with pdfplumber...")
                text = page.extract_text()
                if text and text.strip():
                    extracted_text.append(text)
                    logger.debug(f"Page {page_num + 1}: Extracted {len(text)} characters")
                else:
                    logger.warning(f"Page {page_num + 1}: No extractable text found with pdfplumber")
                    extracted_text.append(f"Page {page_num + 1}: No extractable text found")
            # Check if we got meaningful text
            total_chars = sum(len(text) for text in extracted_text if not text.startswith("Page"))
            if total_chars > 100:  # If we got substantial text
                logger.info(f"✅ pdfplumber extraction successful: {total_chars} characters from {len(extracted_text)} pages")
                logger.info(f"⏱️ PDF extraction completed in {time.time() - start_time:.2f}s")
                return extracted_text
            else:
                logger.warning("pdfplumber extracted minimal text, trying fallback methods...")
    except Exception as pdf_error:
        last_error = pdf_error
        logger.error(f"❌ pdfplumber failed: {pdf_error}")

    # Fallback to pymupdf
    logger.info("📖 Attempting text extraction with pymupdf...")
    try:
        with pymupdf.open(stream=pdf_bytes, filetype="pdf") as doc:
            extracted_text = []
            for page_num in range(len(doc)):
                logger.debug(f"Processing page {page_num + 1} with pymupdf...")
                page = doc[page_num]
                text = page.get_text()
                if text and text.strip():
                    extracted_text.append(text)
                    logger.debug(f"Page {page_num + 1}: Extracted {len(text)} characters")
                else:
                    logger.warning(f"Page {page_num + 1}: No extractable text found with pymupdf")
                    extracted_text.append(f"Page {page_num + 1}: No extractable text found")
            # Check if we got meaningful text
            total_chars = sum(len(text) for text in extracted_text if not text.startswith("Page"))
            if total_chars > 100:  # If we got substantial text
                logger.info(f"✅ pymupdf extraction successful: {total_chars} characters from {len(extracted_text)} pages")
                logger.info(f"⏱️ PDF extraction completed in {time.time() - start_time:.2f}s")
                return extracted_text
            else:
                logger.warning("pymupdf extracted minimal text, trying OCR...")
    except Exception as pymupdf_error:
        last_error = pymupdf_error
        logger.error(f"❌ pymupdf failed: {pymupdf_error}")

    # Last resort: OCR with pytesseract
    logger.info("🔍 Attempting OCR extraction with pytesseract...")
    try:
        # Convert PDF to images
        logger.info("Converting PDF to images for OCR...")
        images = convert_from_bytes(pdf_bytes, dpi=200, fmt='jpeg')
        logger.info(f"Converted PDF to {len(images)} images")
        extracted_text = []
        for page_num, image in enumerate(images):
            logger.debug(f"Running OCR on page {page_num + 1}...")
            try:
                # Run OCR on the image
                text = pytesseract.image_to_string(image, lang='eng')
                if text and text.strip():
                    extracted_text.append(text)
                    logger.debug(f"Page {page_num + 1}: OCR extracted {len(text)} characters")
                else:
                    logger.warning(f"Page {page_num + 1}: No text found via OCR")
                    extracted_text.append(f"Page {page_num + 1}: No text found via OCR")
            except Exception as ocr_page_error:
                last_error = ocr_page_error
                logger.error(f"❌ OCR failed for page {page_num + 1}: {ocr_page_error}")
                extracted_text.append(f"Page {page_num + 1}: OCR processing failed")
        total_chars = sum(len(text) for text in extracted_text if not text.startswith("Page"))
        logger.info(f"✅ OCR extraction completed: {total_chars} characters from {len(extracted_text)} pages")
        logger.info(f"⏱️ PDF extraction completed in {time.time() - start_time:.2f}s")
        return extracted_text
    except Exception as ocr_error:
        last_error = ocr_error
        logger.error(f"❌ OCR extraction failed: {ocr_error}")

    # If all methods fail
    logger.error("❌ All PDF text extraction methods failed")
    logger.info(f"⏱️ PDF extraction failed after {time.time() - start_time:.2f}s")
    return [f"Error: All text extraction methods failed. Last error: {last_error}"]

def create_records(page_data: str, system_prompt: str):
    try:
        datarecords = generation_agent(page_data, system_prompt=system_prompt)
        return datarecords
    except Exception as e:
        print(f"QA generation failed for a page: {str(e)}")
    return []

async def generate_full_dataset_async(chunks: list, schema: list, num_records: int, task_id: str = None):
    """
    Generate dataset records based on document chunks and schema definition using parallel processing.
    
    Args:
        chunks: List of document chunks with content
        schema: Schema definition with field specifications
        num_records: Number of records to generate
        task_id: Optional task ID for progress updates
    
    Returns:
        List of generated dataset records
    """
    logger.info(f"🎯 Starting async dataset generation: {num_records} records from {len(chunks)} chunks")
    
    # Debug: Log detailed chunk information to diagnose count discrepancy
    logger.info(f"📊 Chunk Analysis:")
    logger.info(f"   Total chunks received: {len(chunks)}")
    
    # Analyze chunk sources
    chunk_sources = {}
    for i, chunk in enumerate(chunks):
        filename = chunk.get('filename', 'unknown')
        chunk_sources[filename] = chunk_sources.get(filename, 0) + 1
        logger.debug(f"   Chunk {i+1}: {filename} (type: {chunk.get('chunk_type', 'unknown')})")
    
    for filename, count in chunk_sources.items():
        logger.info(f"   📄 {filename}: {count} chunks")
    
    logger.info(f"   📈 Records/Chunks ratio: {num_records}/{len(chunks)} = {num_records/len(chunks):.2f}")
    
    if num_records < len(chunks):
        logger.warning(f"⚠️  INEFFICIENT: Requesting {num_records} records from {len(chunks)} chunks")
        logger.warning(f"   This will waste {len(chunks) - num_records} chunks ({((len(chunks) - num_records) / len(chunks)) * 100:.1f}% of content)")
    elif num_records > len(chunks):
        logger.info(f"📈 Records > Chunks: Will generate {num_records // len(chunks)} base + {num_records % len(chunks)} extra records")
    
    if not chunks:
        logger.error("No chunks provided for dataset generation")
        return []
    
    if not schema:
        logger.error("No schema provided for dataset generation")
        return []
    
    # Create schema description (same approach as sample preview)
    schema_fields_desc = []
    for f in schema:
        schema_fields_desc.append(f"- {f.get('key','')}: {f.get('description','')} (type: {f.get('type','string')})")
    schema_description = "\n".join(schema_fields_desc)
    
    # Use the same system prompt approach as sample preview (line 707 in document_routes.py)
    system_prompt = f"""You are a dataset generation expert. Generate structured data records based on the provided document content.

**Required Schema:**
{schema_description}

**Instructions:**
1. Analyze the provided document content carefully
2. Generate realistic records that follow the exact schema above
3. Each record must include ALL schema fields
4. Values should be relevant to the document content
5. Ensure data quality and consistency
6. Return valid JSON array format

**Output Format:**
Return a JSON array where each object represents one record with all schema fields."""

    # Process chunks with smart distribution
    # Smart distribution: Handle cases where records < chunks efficiently
    if num_records < len(chunks):
        logger.info(f"🧠 Smart Distribution: {num_records} records < {len(chunks)} chunks")
        
        # Strategy 1: Select chunks with most content (quality-based selection)
        chunks_with_size = []
        for i, chunk in enumerate(chunks):
            content = chunk.get('page_content', '') or chunk.get('content', '')
            content_length = len(content.strip()) if content else 0
            chunks_with_size.append((i, chunk, content_length))
        
        # Sort by content length (descending) and select top N chunks
        chunks_with_size.sort(key=lambda x: x[2], reverse=True)
        selected_chunks = chunks_with_size[:num_records]
        
        logger.info(f"   📊 Selected top {num_records} chunks by content length:")
        for i, (orig_idx, chunk, length) in enumerate(selected_chunks):
            filename = chunk.get('filename', 'unknown')
            logger.info(f"   📄 Rank {i+1}: {filename} (original index {orig_idx+1}, {length} chars)")
        
        # Use selected chunks with 1 record each
        working_chunks = [chunk for _, chunk, _ in selected_chunks]
        chunk_record_counts = [1] * len(working_chunks)
        
        logger.info(f"   📈 Distribution plan: {chunk_record_counts} (total: {sum(chunk_record_counts)} records)")
        
    else:
        # Original balanced distribution for records >= chunks
        logger.info(f"📊 Balanced Distribution: {num_records} records >= {len(chunks)} chunks")
        
        working_chunks = chunks
        base_records_per_chunk = num_records // len(chunks)
        remainder = num_records % len(chunks)
        
        logger.info(f"   Base records per chunk: {base_records_per_chunk}")
        logger.info(f"   Remainder to distribute: {remainder}")
        
        # Create distribution plan
        chunk_record_counts = []
        for i in range(len(chunks)):
            # First 'remainder' chunks get one extra record
            records_for_chunk = base_records_per_chunk + (1 if i < remainder else 0)
            chunk_record_counts.append(records_for_chunk)
        
        logger.info(f"   📈 Distribution plan: {chunk_record_counts} (total: {sum(chunk_record_counts)} records)")
    
    # Create async tasks for parallel processing
    async def process_chunk(i, chunk, batch_records):
        if batch_records <= 0:
            logger.info(f"Processing chunk {i+1}/{len(working_chunks)}: skipping (0 records assigned)")
            return []
            
        logger.info(f"Processing chunk {i+1}/{len(working_chunks)}: generating {batch_records} records")
        
        # Prepare content for generation
        chunk_content = chunk.get('page_content', '') or chunk.get('content', '')
        if not chunk_content:
            logger.warning(f"Chunk {i+1} has no content, skipping...")
            return []
            
        # Truncate content if too long (OpenAI token limits)
        max_content_length = 4000  # Reduced limit to prevent API issues
        if len(chunk_content) > max_content_length:
            chunk_content = chunk_content[:max_content_length] + "..."
            
        generation_prompt = f"""Document Content:
{chunk_content}

Generate {batch_records} records based on this content following the schema above."""

        try:
            # Generate records using the async generation agent
            logger.info(f"Calling async generation_agent for chunk {i+1}...")
            batch_result = await generation_agent_async(
                content=generation_prompt,
                system_prompt=system_prompt,
                model="gpt-5-nano"
                # model="gpt-5-nano"
            )
            logger.info(f"Async generation agent returned: {type(batch_result)} - {str(batch_result)[:100]}...")
            
            # Parse the result
            if batch_result:
                if isinstance(batch_result, str):
                    try:
                        batch_records_data = json.loads(batch_result)
                        if isinstance(batch_records_data, list):
                            actual_added = len(batch_records_data[:batch_records])
                            logger.info(f"Generated {actual_added} records from chunk {i+1}")
                            return batch_records_data[:batch_records]
                        else:
                            logger.warning(f"Chunk {i+1}: Expected list, got {type(batch_records_data)}")
                            return []
                    except json.JSONDecodeError as e:
                        logger.error(f"Chunk {i+1}: JSON parsing failed: {e}")
                        logger.error(f"Raw response: {batch_result[:200]}...")
                        return []
                elif isinstance(batch_result, list):
                    actual_added = len(batch_result[:batch_records])
                    logger.info(f"Generated {actual_added} records from chunk {i+1}")
                    return batch_result[:batch_records]
            return []
                    
        except Exception as e:
            logger.error(f"Failed to generate records from chunk {i+1}: {e}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            return []
    
    # Execute all chunk processing tasks in parallel
    logger.info(f"🚀 Starting parallel processing of {len(working_chunks)} selected chunks...")
    start_time = time.time()
    
    tasks = []
    for i, chunk in enumerate(working_chunks):
        batch_records = chunk_record_counts[i]
        task = process_chunk(i, chunk, batch_records)
        tasks.append(task)
    
    # Use asyncio.gather for parallel execution
    chunk_results = await asyncio.gather(*tasks, return_exceptions=True)
    
    # Process results and handle exceptions
    dataset = []
    for i, result in enumerate(chunk_results):
        if isinstance(result, Exception):
            logger.error(f"Chunk {i+1} failed with exception: {result}")
        elif isinstance(result, list):
            dataset.extend(result)
        else:
            logger.warning(f"Chunk {i+1} returned unexpected type: {type(result)}")
    
    parallel_time = time.time() - start_time
    logger.info(f"⚡ Parallel processing completed in {parallel_time:.2f}s")
    
    # If we still need more records and have some data, create simple variations
    if len(dataset) < num_records and len(dataset) > 0:
        logger.info(f"Generated {len(dataset)}/{num_records} records. Creating variations...")
        
        # Create variations up to the limit (with safety check)
        attempts = 0
        max_attempts = num_records * 2  # Safety limit to prevent infinite loop
        
        while len(dataset) < num_records and attempts < max_attempts:
            attempts += 1
            try:
                import random
                base_record = random.choice(dataset)
                new_record = base_record.copy()
                
                # Add slight variation to avoid exact duplicates
                for field in schema:
                    field_key = field.get('key', '')
                    if field_key in new_record and field.get('type') == 'string':
                        original_value = str(new_record[field_key])
                        new_record[field_key] = f"{original_value} (variation {attempts})"
                        break
                
                dataset.append(new_record)
            except Exception as e:
                logger.error(f"Error creating variation: {e}")
                break
    
    # Trim to exact number if we generated too many
    dataset = dataset[:num_records]
    
    logger.info(f"✅ Async dataset generation completed: {len(dataset)} records generated in {parallel_time:.2f}s")
    
    return dataset

def generate_full_dataset(chunks: list, schema: list, num_records: int, task_id: str = None):
    """
    Synchronous wrapper for generate_full_dataset_async for backwards compatibility.
    Handles different async contexts properly.
    
    Args:
        chunks: List of document chunks with content
        schema: Schema definition with field specifications
        num_records: Number of records to generate
        task_id: Optional task ID for progress updates
    
    Returns:
        List of generated dataset records
    """
    try:
        # Check if we're already in an async context
        try:
            loop = asyncio.get_running_loop()
            logger.info("🔄 Detected existing event loop - using asyncio.create_task()")
            
            # We're in an async context, but this function is called synchronously
            # This is a problematic case - we should recommend using the async version directly
            logger.warning("⚠️  Called sync wrapper from async context - this may cause issues")
            logger.warning("   Recommend calling generate_full_dataset_async() directly from async routes")
            
            # For now, fall back to sequential to avoid conflicts
            logger.info("🔽 Falling back to sequential processing to avoid event loop conflicts")
            return generate_full_dataset_sequential(chunks, schema, num_records, task_id)
            
        except RuntimeError:
            # No event loop running - we can create our own
            logger.info("🆕 No existing event loop - creating new one for async processing")
            
            # Use asyncio.run() which properly handles loop lifecycle
            return asyncio.run(generate_full_dataset_async(chunks, schema, num_records, task_id))
            
    except Exception as e:
        logger.error(f"Error in sync wrapper: {e}")
        # Fallback to original sequential implementation
        logger.warning("🔽 Falling back to sequential processing due to error...")
        return generate_full_dataset_sequential(chunks, schema, num_records, task_id)

def generate_full_dataset_sequential(chunks: list, schema: list, num_records: int, task_id: str = None):
    """
    Sequential implementation with smart distribution as fallback.
    """
    logger.info(f"🎯 Starting sequential dataset generation: {num_records} records from {len(chunks)} chunks")
    
    # Add the same chunk analysis as async version
    logger.info(f"📊 Chunk Analysis:")
    logger.info(f"   Total chunks received: {len(chunks)}")
    
    # Analyze chunk sources
    chunk_sources = {}
    for i, chunk in enumerate(chunks):
        filename = chunk.get('filename', 'unknown')
        chunk_sources[filename] = chunk_sources.get(filename, 0) + 1
        logger.debug(f"   Chunk {i+1}: {filename} (type: {chunk.get('chunk_type', 'unknown')})")
    
    for filename, count in chunk_sources.items():
        logger.info(f"   📄 {filename}: {count} chunks")
    
    logger.info(f"   📈 Records/Chunks ratio: {num_records}/{len(chunks)} = {num_records/len(chunks):.2f}")
    
    if num_records < len(chunks):
        logger.warning(f"⚠️  INEFFICIENT: Requesting {num_records} records from {len(chunks)} chunks")
        logger.warning(f"   This will waste {len(chunks) - num_records} chunks ({((len(chunks) - num_records) / len(chunks)) * 100:.1f}% of content)")
    elif num_records > len(chunks):
        logger.info(f"📈 Records > Chunks: Will generate {num_records // len(chunks)} base + {num_records % len(chunks)} extra records")
    
    if not chunks:
        logger.error("No chunks provided for dataset generation")
        return []
    
    if not schema:
        logger.error("No schema provided for dataset generation")
        return []
    
    dataset = []
    
    # Create schema description (same approach as sample preview)
    schema_fields_desc = []
    for f in schema:
        schema_fields_desc.append(f"- {f.get('key','')}: {f.get('description','')} (type: {f.get('type','string')})")
    schema_description = "\n".join(schema_fields_desc)
    
    # Use the same system prompt approach as sample preview (line 707 in document_routes.py)
    system_prompt = f"""You are a dataset generation expert. Generate structured data records based on the provided document content.

**Required Schema:**
{schema_description}

**Instructions:**
1. Analyze the provided document content carefully
2. Generate realistic records that follow the exact schema above
3. Each record must include ALL schema fields
4. Values should be relevant to the document content
5. Ensure data quality and consistency
6. Return valid JSON array format

**Output Format:**
Return a JSON array where each object represents one record with all schema fields."""

    # Use SMART distribution (same logic as async version)
    if num_records < len(chunks):
        logger.info(f"🧠 Smart Distribution: {num_records} records < {len(chunks)} chunks")
        
        # Strategy 1: Select chunks with most content (quality-based selection)
        chunks_with_size = []
        for i, chunk in enumerate(chunks):
            content = chunk.get('page_content', '') or chunk.get('content', '')
            content_length = len(content.strip()) if content else 0
            chunks_with_size.append((i, chunk, content_length))
        
        # Sort by content length (descending) and select top N chunks
        chunks_with_size.sort(key=lambda x: x[2], reverse=True)
        selected_chunks = chunks_with_size[:num_records]
        
        logger.info(f"   📊 Selected top {num_records} chunks by content length:")
        for i, (orig_idx, chunk, length) in enumerate(selected_chunks):
            filename = chunk.get('filename', 'unknown')
            logger.info(f"   📄 Rank {i+1}: {filename} (original index {orig_idx+1}, {length} chars)")
        
        # Use selected chunks with 1 record each
        working_chunks = [chunk for _, chunk, _ in selected_chunks]
        chunk_record_counts = [1] * len(working_chunks)
        
        logger.info(f"   📈 Distribution plan: {chunk_record_counts} (total: {sum(chunk_record_counts)} records)")
        
    else:
        # Original balanced distribution for records >= chunks
        logger.info(f"📊 Balanced Distribution: {num_records} records >= {len(chunks)} chunks")
        
        working_chunks = chunks
        base_records_per_chunk = num_records // len(chunks)
        remainder = num_records % len(chunks)
        
        logger.info(f"   Base records per chunk: {base_records_per_chunk}")
        logger.info(f"   Remainder to distribute: {remainder}")
        
        # Create distribution plan
        chunk_record_counts = []
        for i in range(len(chunks)):
            # First 'remainder' chunks get one extra record
            records_for_chunk = base_records_per_chunk + (1 if i < remainder else 0)
            chunk_record_counts.append(records_for_chunk)
        
        logger.info(f"   📈 Distribution plan: {chunk_record_counts} (total: {sum(chunk_record_counts)} records)")
    
    for i, chunk in enumerate(working_chunks):
        batch_records = chunk_record_counts[i]
        
        if batch_records <= 0:
            logger.info(f"Processing chunk {i+1}/{len(working_chunks)}: skipping (0 records assigned)")
            continue
            
        logger.info(f"Processing chunk {i+1}/{len(working_chunks)}: generating {batch_records} records")
        
        # Prepare content for generation
        chunk_content = chunk.get('page_content', '') or chunk.get('content', '')
        if not chunk_content:
            logger.warning(f"Chunk {i+1} has no content, skipping...")
            continue
            
        # Truncate content if too long (OpenAI token limits)
        max_content_length = 4000  # Reduced limit to prevent API issues
        if len(chunk_content) > max_content_length:
            chunk_content = chunk_content[:max_content_length] + "..."
            
        generation_prompt = f"""Document Content:
{chunk_content}

Generate {batch_records} records based on this content following the schema above."""

        try:
            # Generate records using the generation agent
            logger.info(f"Calling generation_agent for chunk {i+1}...")
            batch_result = generation_agent(
                content=generation_prompt,
                system_prompt=system_prompt,
                model="gpt-5-nano"
                # model="gpt-5-nano"
            )
            logger.info(f"Generation agent returned: {type(batch_result)} - {str(batch_result)[:100]}...")
            
            # Parse the result
            if batch_result:
                if isinstance(batch_result, str):
                    try:
                        batch_records_data = json.loads(batch_result)
                        if isinstance(batch_records_data, list):
                            actual_added = len(batch_records_data[:batch_records])
                            dataset.extend(batch_records_data[:batch_records])
                            logger.info(f"Generated {actual_added} records from chunk {i+1}")
                        else:
                            logger.warning(f"Chunk {i+1}: Expected list, got {type(batch_records_data)}")
                    except json.JSONDecodeError as e:
                        logger.error(f"Chunk {i+1}: JSON parsing failed: {e}")
                        logger.error(f"Raw response: {batch_result[:200]}...")
                elif isinstance(batch_result, list):
                    actual_added = len(batch_result[:batch_records])
                    dataset.extend(batch_result[:batch_records])
                    logger.info(f"Generated {actual_added} records from chunk {i+1}")
                    
        except Exception as e:
            logger.error(f"Failed to generate records from chunk {i+1}: {e}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            continue
    
    # If we still need more records and have some data, create simple variations
    if len(dataset) < num_records and len(dataset) > 0:
        logger.info(f"Generated {len(dataset)}/{num_records} records. Creating variations...")
        
        # Create variations up to the limit (with safety check)
        attempts = 0
        max_attempts = num_records * 2  # Safety limit to prevent infinite loop
        
        while len(dataset) < num_records and attempts < max_attempts:
            attempts += 1
            try:
                import random
                base_record = random.choice(dataset)
                new_record = base_record.copy()
                
                # Add slight variation to avoid exact duplicates
                for field in schema:
                    field_key = field.get('key', '')
                    if field_key in new_record and field.get('type') == 'string':
                        original_value = str(new_record[field_key])
                        new_record[field_key] = f"{original_value} (variation {attempts})"
                        break
                
                dataset.append(new_record)
            except Exception as e:
                logger.error(f"Error creating variation: {e}")
                break
    
    # Trim to exact number if we generated too many
    dataset = dataset[:num_records]
    
    logger.info(f"✅ Sequential dataset generation completed: {len(dataset)} records generated")
    
    return dataset